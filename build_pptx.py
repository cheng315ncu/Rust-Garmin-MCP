"""Build PRESENTATION.pptx — refined Tesla-style deck for Garmin MCP (Rust).

Unified light palette. Editorial composition. SF Pro Display / SF Pro Text /
SF Mono for distinctive Apple-industrial feel. Single Electric-Blue accent.
Top-edge blue progress bar + bottom tracker as a deck-wide signature.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ───── Tesla palette (light only) ─────
BLUE   = RGBColor(0x3E, 0x6A, 0xE1)
BLUE_D = RGBColor(0x32, 0x57, 0xB8)
CARBON = RGBColor(0x17, 0x1A, 0x20)
GRAPH  = RGBColor(0x39, 0x3C, 0x41)
PEWTER = RGBColor(0x5C, 0x5E, 0x62)
SILVER = RGBColor(0x8E, 0x8E, 0x8E)
CLOUD  = RGBColor(0xEE, 0xEE, 0xEE)
PALE   = RGBColor(0xD0, 0xD1, 0xD2)
ASH    = RGBColor(0xF6, 0xF6, 0xF7)   # very subtle warm white
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

# Distinctive Apple-native typography (every macOS has these, not generic AI fonts)
DISPLAY = "SF Pro Display"
TEXT    = "SF Pro Text"
MONO    = "SF Mono"

TOTAL_SLIDES = 9

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


# ──── shape primitives ────
def kill_shadow(shape):
    spPr = shape.fill._xPr
    existing = spPr.find(qn('a:effectLst'))
    if existing is not None:
        spPr.remove(existing)
    etree.SubElement(spPr, qn('a:effectLst'))


def add_slide(bg=WHITE, slide_no=1):
    s = prs.slides.add_slide(blank)
    bg_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg_shape.line.fill.background()
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    kill_shadow(bg_shape)

    # Top-edge blue progress bar — deck-wide signature, 2pt tall
    progress = bg_shape  # ref
    pb_w = SW * (slide_no / TOTAL_SLIDES)
    pb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, pb_w, Emu(20000))
    pb.line.fill.background()
    pb.fill.solid(); pb.fill.fore_color.rgb = BLUE
    kill_shadow(pb)

    # Subtle background-tint progress trail (rest of bar, very pale)
    pb_rest = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, pb_w, 0, SW - pb_w, Emu(20000))
    pb_rest.line.fill.background()
    pb_rest.fill.solid(); pb_rest.fill.fore_color.rgb = CLOUD
    kill_shadow(pb_rest)

    return s


def add_text(slide, x, y, w, h, text, *, font=TEXT, size=14, bold=False,
             italic=False, color=CARBON, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, tracking=0, leading=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if leading:
            p.line_spacing = leading
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if tracking:
            rPr = run._r.get_or_add_rPr()
            rPr.set('spc', str(tracking))
    return tb


def add_rect(slide, x, y, w, h, *, fill=WHITE, line=None, line_w=0.5,
             radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    if radius is not None:
        shp.adjustments[0] = radius
    kill_shadow(shp)
    return shp


def add_line(slide, x1, y1, x2, y2, color=CLOUD, weight=0.5):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


# ──── deck-wide chrome ────
def add_chrome(slide, slide_no):
    """Bottom slide tracker only — top header removed for more vertical room."""
    # Bottom tracker — 9 tiny dashes, current = blue
    bx, by = 11.05, 7.15
    bw, bh, bg = 0.16, 0.025, 0.04
    for i in range(TOTAL_SLIDES):
        add_rect(slide, Inches(bx + i * (bw + bg)), Inches(by),
                 Inches(bw), Inches(bh),
                 fill=BLUE if i + 1 == slide_no else PALE)
    # Slide number — monospace
    add_text(slide, Inches(0.7), Inches(7.10),
             Inches(2), Inches(0.25),
             f"{slide_no:02d} / {TOTAL_SLIDES:02d}",
             font=MONO, size=9, color=PEWTER, tracking=200)


def add_eyebrow(slide, x, y, num, label):
    """Editorial eyebrow: '01 — THE MANDATE' (mono number + tracked label)."""
    add_text(slide, x, y, Inches(8), Inches(0.3),
             f"{num:02d}  —  {label.upper()}",
             font=MONO, size=10, bold=True, color=BLUE, tracking=400)


def cta_button(slide, x, y, label, *, primary=True, w=2.0, h=0.5):
    fill = BLUE if primary else WHITE
    line = None if primary else PALE
    txt_color = WHITE if primary else GRAPH
    add_rect(slide, Inches(x), Inches(y), Inches(w), Inches(h),
             fill=fill, line=line, radius=0.18)
    add_text(slide, Inches(x), Inches(y), Inches(w), Inches(h),
             label, font=TEXT, size=12, bold=True, color=txt_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def code_block(slide, x, y, w, h, lines, size=10.5):
    """Light code block: ash bg, carbon text, blue keywords, pewter comments."""
    add_rect(slide, x, y, w, h, fill=ASH,
             line=CLOUD, line_w=0.5, radius=0.04)
    tb = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.22),
                                   w - Inches(0.6), h - Inches(0.44))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.45
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = MONO
        run.font.size = Pt(size)
        if line.lstrip().startswith("//"):
            run.font.color.rgb = PEWTER
            run.font.italic = True
        elif line.lstrip().startswith(("pub ", "let ", "fn ", "use ", "struct ", "impl ")):
            run.font.color.rgb = BLUE
            run.font.bold = True
        else:
            run.font.color.rgb = CARBON


# ═════════════════════════════════════════════════════════════
# SLIDE 1 — HERO  (white)  · editorial center composition
# ═════════════════════════════════════════════════════════════
s = add_slide(WHITE, 1)
add_chrome(s, 1)

# Top tag — small, centered, mono for editorial cool
add_text(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.4),
         "MODEL CONTEXT PROTOCOL  ·  GARMIN CONNECT",
         font=MONO, size=10, bold=True, color=PEWTER,
         align=PP_ALIGN.CENTER, tracking=400)

# Display title — SF Pro Display Heavy
add_text(s, Inches(0.5), Inches(1.50), Inches(12.3), Inches(1.2),
         "Garmin MCP",
         font=DISPLAY, size=78, bold=True, color=CARBON,
         align=PP_ALIGN.CENTER)

# Blue subhead — italic for editorial accent
add_text(s, Inches(0.5), Inches(2.65), Inches(12.3), Inches(1.0),
         "Engineered in Rust.",
         font=DISPLAY, size=48, italic=True, color=BLUE,
         align=PP_ALIGN.CENTER)

# Lede
add_text(s, Inches(2.5), Inches(3.80), Inches(8.3), Inches(0.8),
         "A single static binary that exposes 77 fitness & health tools to "
         "Claude — with type-safe concurrency, in-process caching, and rate "
         "limiting built in.",
         font=TEXT, size=14, color=GRAPH,
         align=PP_ALIGN.CENTER, leading=1.5)

# Stats row — refined hairlines, tabular display numerals
stats = [("77",     "TOOLS"),
         ("12",     "MODULES"),
         ("~5 MB",  "MEMORY"),
         ("~10 MB", "BINARY")]
sx, sy, sw, sh_ = 2.8, 5.00, 1.93, 0.95
total_w = sw * 4
add_line(s, Inches(sx), Inches(sy),
         Inches(sx + total_w), Inches(sy), color=CARBON, weight=0.75)
add_line(s, Inches(sx), Inches(sy + sh_),
         Inches(sx + total_w), Inches(sy + sh_), color=PALE, weight=0.5)
for i, (v, l) in enumerate(stats):
    cx = sx + i * sw
    add_text(s, Inches(cx), Inches(sy + 0.10),
             Inches(sw), Inches(0.5),
             v, font=DISPLAY, size=26, bold=True, color=CARBON,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(cx), Inches(sy + 0.62),
             Inches(sw), Inches(0.3),
             l, font=MONO, size=8.5, bold=True, color=PEWTER,
             align=PP_ALIGN.CENTER, tracking=400)

# CTAs — centered
cta_button(s, 4.6, 6.30, "View Architecture", primary=True, w=2.0, h=0.45)
cta_button(s, 6.75, 6.30, "Why Rust", primary=False, w=2.0, h=0.45)


# ═════════════════════════════════════════════════════════════
# SLIDE 2 — WHY RUST  (white)
# ═════════════════════════════════════════════════════════════
s = add_slide(WHITE, 2)
add_chrome(s, 2)
add_eyebrow(s, Inches(0.7), Inches(0.55), 1, "The Mandate")

# Headline left
add_text(s, Inches(0.7), Inches(1.10), Inches(5.4), Inches(2.4),
         ["One binary.", "Zero runtime.", "Predictable at 3 a.m."],
         font=DISPLAY, size=38, bold=True, color=CARBON, leading=1.08)

add_text(s, Inches(0.7), Inches(3.90), Inches(5.2), Inches(2.0),
         "A Rust port of the Python original. Built so an LLM can hammer "
         "Garmin Connect without triggering a session lockout, leaking "
         "credentials, or paging you to restart a Python venv.",
         font=TEXT, size=12.5, color=GRAPH, leading=1.6)

# Comparison table right — refined
rows = [
    ("Deployment",        "Python 3.12 + uv/pip",   "Single static binary"),
    ("Memory",            "~50 MB",                 "~5 MB"),
    ("Binary size",       "N/A",                    "~10 MB release"),
    ("Session safety",    "GIL + race conditions",  "Mutex / RwLock at type"),
    ("Duplicate queries", "Each tool re-fetches",   "moka TTL + singleflight"),
    ("Rate limiting",     "Burst → lockout",        "governor token-bucket"),
    ("Research output",   "JSON only",              "JSON · CSV · EDF-ready"),
]
tx, ty, tw = 6.5, 1.05, 6.15
col_w = [1.95, 2.10, 2.10]
header_h = 0.42
row_h = 0.55

# Card behind table
card_h = header_h + len(rows) * row_h + 0.45
add_rect(s, Inches(tx - 0.25), Inches(ty - 0.2),
         Inches(tw + 0.5), Inches(card_h),
         fill=ASH, line=CLOUD, radius=0.04)

# Headers — monospace small caps
add_text(s, Inches(tx + col_w[0]), Inches(ty),
         Inches(col_w[1]), Inches(header_h),
         "PYTHON", font=MONO, size=9, bold=True, color=PEWTER,
         tracking=400, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(tx + col_w[0] + col_w[1]), Inches(ty),
         Inches(col_w[2]), Inches(header_h),
         "RUST · THIS", font=MONO, size=9, bold=True, color=BLUE,
         tracking=400, anchor=MSO_ANCHOR.MIDDLE)
add_line(s, Inches(tx), Inches(ty + header_h),
         Inches(tx + tw), Inches(ty + header_h), color=PALE, weight=0.75)

for i, (lbl, py, rs) in enumerate(rows):
    y = ty + header_h + i * row_h
    add_text(s, Inches(tx), Inches(y),
             Inches(col_w[0]), Inches(row_h),
             lbl, font=TEXT, size=11, bold=True, color=CARBON,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(tx + col_w[0]), Inches(y),
             Inches(col_w[1]), Inches(row_h),
             py, font=TEXT, size=10.5, color=SILVER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(tx + col_w[0] + col_w[1]), Inches(y),
             Inches(col_w[2]), Inches(row_h),
             rs, font=TEXT, size=10.5, bold=True, color=BLUE,
             anchor=MSO_ANCHOR.MIDDLE)
    if i < len(rows) - 1:
        add_line(s, Inches(tx), Inches(y + row_h),
                 Inches(tx + tw), Inches(y + row_h), color=CLOUD, weight=0.5)


# ═════════════════════════════════════════════════════════════
# SLIDE 3 — ARCHITECTURE  (ash) · 6-layer pipeline
# ═════════════════════════════════════════════════════════════
s = add_slide(ASH, 3)
add_chrome(s, 3)
add_eyebrow(s, Inches(0.7), Inches(0.55), 2, "Architecture")

add_text(s, Inches(0.7), Inches(0.90), Inches(12), Inches(0.65),
         "The GET pipeline.",
         font=DISPLAY, size=34, bold=True, color=CARBON)
add_text(s, Inches(0.7), Inches(1.45), Inches(12), Inches(0.4),
         "Six layers. Each one earns its place. Cache hits skip everything below them.",
         font=TEXT, size=12, color=GRAPH)

layers = [
    ("L6", "Tool Layer",
     "77 tools · #[tool_router] · GarminMcpServer · Arc<GarminApiClient>", False),
    ("L5", "ClinicalExport",
     "FlatSummary · HrvPayload · TimeseriesArray · EventTable — JSON or CSV", False),
    ("L4", "moka Async Cache",
     "60 s TTL · 1 000 entries · sorted endpoint?k=v key · singleflight via try_get_with", True),
    ("L3", "governor Rate Limiter",
     "Token bucket · 60 req/min · shared reads + writes · until_ready().await", True),
    ("L2", "Rust Sync Layer",
     "Mutex<GarminClient> for GETs · RwLock<BearerToken> for writes · pooled reqwest", False),
    ("L1", "Garmin Connect API",
     "GET via garmin_client · POST/PUT/DELETE via reqwest · invalidate_all() on write", False),
]
ly, lh, gap = 2.10, 0.72, 0.10
lx, lw = 0.7, 11.95
for i, (lvl, name, desc, accent) in enumerate(layers):
    y = ly + i * (lh + gap)
    add_rect(s, Inches(lx), Inches(y), Inches(lw), Inches(lh),
             fill=WHITE, line=BLUE if accent else CLOUD,
             line_w=1.0 if accent else 0.5, radius=0.06)
    # Layer number — monospace blue label
    add_text(s, Inches(lx + 0.3), Inches(y + 0.13),
             Inches(0.8), Inches(0.4),
             lvl, font=MONO, size=11, bold=True,
             color=BLUE if accent else PEWTER, tracking=300)
    # Layer name
    add_text(s, Inches(lx + 1.15), Inches(y + 0.10),
             Inches(2.7), Inches(0.46),
             name, font=DISPLAY, size=15, bold=True,
             color=CARBON, anchor=MSO_ANCHOR.MIDDLE)
    # Description
    add_text(s, Inches(lx + 4.0), Inches(y),
             Inches(lw - 4.2), Inches(lh),
             desc, font=TEXT, size=11, color=GRAPH, leading=1.4,
             anchor=MSO_ANCHOR.MIDDLE)


# ═════════════════════════════════════════════════════════════
# SLIDE 4 — CONCURRENCY  (white)
# ═════════════════════════════════════════════════════════════
s = add_slide(WHITE, 4)
add_chrome(s, 4)
add_eyebrow(s, Inches(0.7), Inches(0.55), 3, "Concurrency Model")

add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.4),
         ["Type-safe sessions.", "By construction."],
         font=DISPLAY, size=38, bold=True, color=CARBON, leading=1.08)

add_text(s, Inches(0.7), Inches(2.40), Inches(11), Inches(0.55),
         "The compiler refuses to ship a second OAuth session. Three primitives — "
         "each one solves a specific failure mode.",
         font=TEXT, size=12.5, color=GRAPH, leading=1.5)

cards = [
    ("01", "Mutex<GarminClient>",
     "Garmin rate-limits per session token, and api_request needs &mut self. "
     "Serialising every GET through one Mutex is the safety property — not a workaround.",
     True),
    ("02", "RwLock<BearerToken>",
     "Writes hold a read lock concurrently. Token refresh acquires the write lock "
     "with double-checked locking. One token; many concurrent POSTs.",
     False),
    ("03", "reqwest::Client",
     "One pooled HTTP client across every write path. TLS handshakes and TCP "
     "connections are reused — not re-paid on every tool call.",
     False),
]
cy, ch = 3.30, 2.5
cw, cgap = 3.85, 0.2
cx0 = 0.7
for i, (idx, title, body, accent) in enumerate(cards):
    x = cx0 + i * (cw + cgap)
    add_rect(s, Inches(x), Inches(cy), Inches(cw), Inches(ch),
             fill=ASH if accent else WHITE,
             line=BLUE if accent else CLOUD,
             line_w=1.0 if accent else 0.5, radius=0.04)
    # Top index number — monospace
    add_text(s, Inches(x + 0.35), Inches(cy + 0.30),
             Inches(1.0), Inches(0.3),
             idx, font=MONO, size=10, bold=True,
             color=BLUE, tracking=300)
    # Vertical accent rule (only on accent card)
    if accent:
        add_rect(s, Inches(x + 0.35), Inches(cy + 0.55),
                 Inches(0.4), Inches(0.04),
                 fill=BLUE)
    # Title
    add_text(s, Inches(x + 0.35), Inches(cy + 0.78),
             Inches(cw - 0.7), Inches(0.45),
             title, font=DISPLAY, size=15, bold=True, color=CARBON)
    # Body
    add_text(s, Inches(x + 0.35), Inches(cy + 1.30),
             Inches(cw - 0.7), Inches(ch - 1.40),
             body, font=TEXT, size=10.5, color=GRAPH, leading=1.5)

# Code block at bottom
code_block(s, Inches(0.7), Inches(6.05), Inches(11.95), Inches(0.85), [
    "pub struct GarminApiClient {",
    "    inner: Arc<Mutex<GarminClient>>,    http: reqwest::Client,    "
    "token: Arc<RwLock<BearerToken>>,    cache: Cache<…>,    limiter: Arc<Limiter>,",
    "}",
], size=10)


# ═════════════════════════════════════════════════════════════
# SLIDE 5 — CACHE & RATE LIMIT  (ash)
# ═════════════════════════════════════════════════════════════
s = add_slide(ASH, 5)
add_chrome(s, 5)
add_eyebrow(s, Inches(0.7), Inches(0.55), 4, "Cache & Rate Limit")

add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.4),
         ["Coalesce duplicates.", "Respect the upstream."],
         font=DISPLAY, size=38, bold=True, color=CARBON, leading=1.08)

stats4 = [
    ("60s",    "CACHE TTL",
     "Coalesces LLM re-asks. Short enough that fresh wearable data isn't hidden.", True),
    ("1 000",  "CACHE ENTRIES",
     "LRU-evicted past the cap. Values are Arc<Value> for cheap clones.", False),
    ("60/min", "RATE BUDGET",
     "Token bucket via governor. Shared across reads and writes — one budget.", True),
    ("1",      "IN-FLIGHT / KEY",
     "moka's try_get_with deduplicates concurrent callers. 10 tools = 1 fetch.", False),
]
sy, sh4 = 3.20, 2.70
sw4, sgap4 = 2.92, 0.16
sx0 = 0.7
total_w5 = sw4 * 4 + sgap4 * 3
# Top hairline (carbon, structural) + bottom hairline (pale, decorative)
add_line(s, Inches(sx0), Inches(sy),
         Inches(sx0 + total_w5), Inches(sy), color=CARBON, weight=0.75)
add_line(s, Inches(sx0), Inches(sy + sh4),
         Inches(sx0 + total_w5), Inches(sy + sh4), color=PALE, weight=0.5)
for i, (v, l, x, accent) in enumerate(stats4):
    cx = sx0 + i * (sw4 + sgap4)
    # Big numeral
    add_text(s, Inches(cx), Inches(sy + 0.25),
             Inches(sw4), Inches(0.85),
             v, font=DISPLAY, size=44, bold=True,
             color=BLUE if accent else CARBON)
    # Label — monospace
    add_text(s, Inches(cx), Inches(sy + 1.18),
             Inches(sw4), Inches(0.3),
             l, font=MONO, size=9, bold=True, color=PEWTER, tracking=400)
    # Description
    add_text(s, Inches(cx), Inches(sy + 1.55),
             Inches(sw4 - 0.15), Inches(sh4 - 1.55),
             x, font=TEXT, size=10.5, color=GRAPH, leading=1.5)
    # Vertical hairline between cells (subtle)
    if i > 0:
        add_line(s, Inches(cx - sgap4 / 2), Inches(sy + 0.25),
                 Inches(cx - sgap4 / 2), Inches(sy + sh4 - 0.25),
                 color=CLOUD, weight=0.5)

# Two-up rationale — refined hairline-led
ry = 6.30
add_line(s, Inches(0.7), Inches(ry - 0.05),
         Inches(6.4), Inches(ry - 0.05), color=BLUE, weight=1.0)
add_text(s, Inches(0.7), Inches(ry), Inches(5.7), Inches(0.32),
         "Why TTL caching",
         font=MONO, size=10, bold=True, color=BLUE, tracking=300)
add_text(s, Inches(0.7), Inches(ry + 0.32), Inches(5.7), Inches(0.55),
         "An LLM asks the same question twice. moka turns repeats into Arc-bumps "
         "— Mutex and network never see them.",
         font=TEXT, size=10.5, color=GRAPH, leading=1.5)

add_line(s, Inches(7.0), Inches(ry - 0.05),
         Inches(12.65), Inches(ry - 0.05), color=BLUE, weight=1.0)
add_text(s, Inches(7.0), Inches(ry), Inches(5.65), Inches(0.32),
         "Why rate limit",
         font=MONO, size=10, bold=True, color=BLUE, tracking=300)
add_text(s, Inches(7.0), Inches(ry + 0.32), Inches(5.65), Inches(0.55),
         "Garmin's per-session lockout is real but undocumented. governor "
         "enforces a conservative ceiling.",
         font=TEXT, size=10.5, color=GRAPH, leading=1.5)


# ═════════════════════════════════════════════════════════════
# SLIDE 6 — SECURITY  (white) · 3×2 refined cards
# ═════════════════════════════════════════════════════════════
s = add_slide(WHITE, 6)
add_chrome(s, 6)
add_eyebrow(s, Inches(0.7), Inches(0.55), 5, "Security Advantages")

add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.4),
         ["Memory safety", "is a security feature."],
         font=DISPLAY, size=38, bold=True, color=CARBON, leading=1.08)

add_text(s, Inches(0.7), Inches(2.40), Inches(11.5), Inches(0.55),
         "No GC pauses leaking timing. No eval. No dynamic dispatch on user "
         "input. No Python supply chain at runtime.",
         font=TEXT, size=12.5, color=GRAPH, leading=1.5)

sec_cards = [
    ("Credentials at rest",
     "Env vars or _FILE mode reading chmod-600 files. Secrets never in client config."),
    ("Type-safe sessions",
     "No data race can produce a second OAuth session — borrow checker rejects at compile."),
    ("Local stdio only",
     "MCP runs over stdin/stdout. No exposed port, no remote attack surface."),
    ("Destructive ops omitted",
     "delete_activity not exposed. get_activity_details (50–500 KB GPS) also omitted."),
    ("Minimal supply chain",
     "11 direct dependencies in Cargo.toml. Reviewable. No pip transitive surprises."),
    ("Rate-limit hardening",
     "A model that loops accidentally cannot brick your account — limiter caps first."),
]
gx, gy = 0.7, 3.10
gw, ggap_x = 3.97, 0.13
ghh, ggap_y = 1.85, 0.18
for idx, (t, d) in enumerate(sec_cards):
    col = idx % 3
    row = idx // 3
    x = gx + col * (gw + ggap_x)
    y = gy + row * (ghh + ggap_y)
    add_rect(s, Inches(x), Inches(y), Inches(gw), Inches(ghh),
             fill=WHITE, line=CLOUD, line_w=0.5, radius=0.04)
    # Top-edge blue accent rule (4px tall, only across left half)
    add_rect(s, Inches(x), Inches(y), Inches(0.5), Inches(0.04),
             fill=BLUE)
    # Number — monospace small
    add_text(s, Inches(x + 0.35), Inches(y + 0.25),
             Inches(1.0), Inches(0.3),
             f"{idx + 1:02d}",
             font=MONO, size=9.5, bold=True, color=BLUE, tracking=300)
    # Title
    add_text(s, Inches(x + 0.35), Inches(y + 0.55),
             Inches(gw - 0.7), Inches(0.45),
             t, font=DISPLAY, size=14, bold=True, color=CARBON)
    # Body
    add_text(s, Inches(x + 0.35), Inches(y + 1.00),
             Inches(gw - 0.7), Inches(ghh - 1.10),
             d, font=TEXT, size=10.5, color=GRAPH, leading=1.5)


# ═════════════════════════════════════════════════════════════
# SLIDE 7 — TOOLS  (ash) · 4×3 module grid
# ═════════════════════════════════════════════════════════════
s = add_slide(ASH, 7)
add_chrome(s, 7)
add_eyebrow(s, Inches(0.7), Inches(0.55), 6, "Tool Coverage")

add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.4),
         ["77 tools.", "12 focused modules."],
         font=DISPLAY, size=38, bold=True, color=CARBON, leading=1.08)

add_text(s, Inches(0.7), Inches(2.40), Inches(11.5), Inches(0.5),
         "Curated, not auto-generated. Each module hides Garmin's endpoint "
         "quirks behind a clean MCP surface.",
         font=TEXT, size=12.5, color=GRAPH, leading=1.5)

modules = [
    ("Health & Wellness", "21",  "stats · sleep · HR · stress · HRV", False),
    ("Activities",        "14",  "by-date · splits · weather · zones", False),
    ("Challenges",         "8",  "badges · ad-hoc · virtual · goals",   False),
    ("Devices",            "6",  "list · primary · settings · alarms",  False),
    ("Workouts",           "5",  "list · get · scheduled · delete",     False),
    ("User Profile",       "4",  "profile · settings · units",          False),
    ("Research",           "4",  "stats · sleep · HRV · weekly",        True),
    ("Training",           "3",  "status · progress · race predictions", False),
    ("Gear",               "3",  "list · attach · detach",              False),
    ("Nutrition",          "3",  "food log · settings · custom foods",  False),
    ("Women's Health",     "3",  "menstrual day · pregnancy",           False),
    ("Data Management",    "3",  "hydration · BP · body composition",   False),
]
mx, my = 0.7, 3.15
mw, mh = 3.0, 1.45
mgap_x, mgap_y = 0.12, 0.16
for i, (n, c, t, accent) in enumerate(modules):
    col = i % 4
    row = i // 4
    x = mx + col * (mw + mgap_x)
    y = my + row * (mh + mgap_y)
    add_rect(s, Inches(x), Inches(y), Inches(mw), Inches(mh),
             fill=WHITE,
             line=BLUE if accent else CLOUD,
             line_w=1.0 if accent else 0.5, radius=0.05)
    # Top-edge accent rule on Research card
    if accent:
        add_rect(s, Inches(x), Inches(y), Inches(0.5), Inches(0.04), fill=BLUE)
    # Module name
    add_text(s, Inches(x + 0.25), Inches(y + 0.22),
             Inches(mw - 1.0), Inches(0.4),
             n, font=DISPLAY, size=12.5, bold=True,
             color=BLUE if accent else CARBON)
    # Big count, right-aligned, blue
    add_text(s, Inches(x + mw - 0.95), Inches(y + 0.15),
             Inches(0.8), Inches(0.55),
             c, font=DISPLAY, size=22, bold=True, color=BLUE,
             align=PP_ALIGN.RIGHT)
    # Tags
    add_text(s, Inches(x + 0.25), Inches(y + 0.72),
             Inches(mw - 0.5), Inches(mh - 0.82),
             t, font=TEXT, size=10, color=PEWTER, leading=1.45)


# ═════════════════════════════════════════════════════════════
# SLIDE 8 — RESEARCH OUTPUT  (white)
# ═════════════════════════════════════════════════════════════
s = add_slide(WHITE, 8)
add_chrome(s, 8)
add_eyebrow(s, Inches(0.7), Inches(0.55), 7, "Research Output")

# Left column
add_text(s, Inches(0.7), Inches(0.95), Inches(5.5), Inches(3.5),
         ["JSON for chat.",
          "CSV for pandas.",
          "EDF on the way."],
         font=DISPLAY, size=36, bold=True, color=CARBON, leading=1.12)

add_text(s, Inches(0.7), Inches(4.05), Inches(5.4), Inches(1.8),
         "A ClinicalExport trait at the output edge means every tool answers "
         "in the format the caller asked for. Up to 366 days of longitudinal "
         "data per call.",
         font=TEXT, size=12.5, color=GRAPH, leading=1.55)

cta_button(s, 0.7, 6.10, "Get Started", primary=True, w=1.85, h=0.5)
cta_button(s, 2.7, 6.10, "Read the Spec", primary=False, w=1.85, h=0.5)

# Right code panel
code_block(s, Inches(6.6), Inches(0.95), Inches(6.05), Inches(6.05), [
    "// 90-day health panel — pandas-ready",
    '{ "tool": "get_daily_stats_range",',
    '  "arguments": {',
    '    "start_date": "2026-01-26",',
    '    "end_date":   "2026-04-25",',
    '    "format":     "csv"',
    '  } }',
    "",
    "// 30-day HRV trends",
    '{ "tool": "get_hrv_range",',
    '  "arguments": {',
    '    "start_date": "2026-03-26",',
    '    "end_date":   "2026-04-25",',
    '    "format":     "csv"',
    '  } }',
    "",
    "// ISO-week aggregates · mean/std/min/max × 12",
    '{ "tool": "get_weekly_summary",',
    '  "arguments": {',
    '    "start_date": "2026-01-01",',
    '    "end_date":   "2026-04-25"',
    '  } }',
])


# ═════════════════════════════════════════════════════════════
# SLIDE 9 — CLOSING  (ash) · final mark
# ═════════════════════════════════════════════════════════════
s = add_slide(ASH, 9)
add_chrome(s, 9)
# Centered eyebrow for the closing slide
add_text(s, Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.3),
         "08  —  READY TO SHIP",
         font=MONO, size=10, bold=True, color=BLUE,
         align=PP_ALIGN.CENTER, tracking=400)

add_text(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.0),
         "One binary.",
         font=DISPLAY, size=58, bold=True, color=CARBON,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(2.30), Inches(12.3), Inches(1.0),
         "One memory model.",
         font=DISPLAY, size=58, bold=True, color=CARBON,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(3.15), Inches(12.3), Inches(1.0),
         "One blue button.",
         font=DISPLAY, size=58, italic=True, bold=True, color=BLUE,
         align=PP_ALIGN.CENTER)

# Decorative blue rule below the headline trio (centered, short)
add_rect(s, Inches(6.17), Inches(4.30), Inches(1.0), Inches(0.025),
         fill=BLUE)

add_text(s, Inches(2.5), Inches(4.50), Inches(8.3), Inches(0.7),
         "Drop the binary into your MCP client config. "
         "No Python. No venv. No surprises at 3 a.m.",
         font=TEXT, size=14, color=GRAPH,
         align=PP_ALIGN.CENTER, leading=1.5)

# CTAs
cta_button(s, 4.6, 5.40, "Order Now", primary=True, w=2.0, h=0.5)
cta_button(s, 6.75, 5.40, "Read the Spec", primary=False, w=2.0, h=0.5)

# Install command — refined mono pill
add_rect(s, Inches(3.5), Inches(6.30), Inches(6.3), Inches(0.5),
         fill=WHITE, line=PALE, radius=0.20)
add_text(s, Inches(3.5), Inches(6.30), Inches(6.3), Inches(0.5),
         "$  cargo install --path .",
         font=MONO, size=11, color=CARBON,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ───── Save ─────
out = "/Users/cheng/alg/mcp/garmin_mcp_rust/PRESENTATION.pptx"
prs.save(out)
print(f"✓ Saved {out}  ({len(prs.slides)} slides)")
